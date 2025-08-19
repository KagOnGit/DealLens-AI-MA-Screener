from fastapi import APIRouter, Depends, HTTPException, Response
from sqlalchemy.orm import Session
from typing import List
from pydantic import BaseModel
from datetime import datetime
import io
from reportlab.pdfgen import canvas
from reportlab.lib.pagesizes import letter
from pptx import Presentation
from pptx.util import Inches

from ....core.deps import get_db

router = APIRouter()


class TearsheetRequest(BaseModel):
    ticker: str
    format: str  # pdf or pptx
    sections: List[str]  # overview, comps, precedents, ownership, news, valuation


@router.post("/tearsheet")
async def generate_tearsheet(
    request: TearsheetRequest,
    db: Session = Depends(get_db)
):
    """Generate and download PDF or PPT tearsheet."""
    try:
        if request.format.lower() == "pdf":
            return generate_pdf_tearsheet(request)
        elif request.format.lower() == "pptx":
            return generate_ppt_tearsheet(request)
        else:
            raise HTTPException(status_code=400, detail="Format must be 'pdf' or 'pptx'")
    except Exception as e:
        raise HTTPException(status_code=500, detail=str(e))


def generate_pdf_tearsheet(request: TearsheetRequest):
    """Generate PDF tearsheet using reportlab."""
    buffer = io.BytesIO()
    p = canvas.Canvas(buffer, pagesize=letter)
    
    # Title page
    width, height = letter
    p.setFont("Helvetica-Bold", 24)
    p.drawString(100, height - 100, f"{request.ticker} - Investment Tearsheet")
    
    p.setFont("Helvetica", 12)
    p.drawString(100, height - 150, f"Generated: {datetime.now().strftime('%B %d, %Y')}")
    p.drawString(100, height - 170, "DealLens - Investment Banking Platform")
    
    # Add sections
    y_position = height - 220
    p.setFont("Helvetica-Bold", 16)
    p.drawString(100, y_position, "Sections Included:")
    
    y_position -= 30
    p.setFont("Helvetica", 12)
    for section in request.sections:
        p.drawString(120, y_position, f"• {section.title()}")
        y_position -= 20
    
    # Mock content for each section
    y_position -= 40
    for section in request.sections:
        if y_position < 100:  # Start new page if needed
            p.showPage()
            y_position = height - 100
        
        p.setFont("Helvetica-Bold", 14)
        p.drawString(100, y_position, f"{section.title()} Analysis")
        y_position -= 25
        
        p.setFont("Helvetica", 10)
        mock_content = f"This section contains detailed {section} analysis for {request.ticker}."
        p.drawString(100, y_position, mock_content)
        y_position -= 60
    
    p.save()
    buffer.seek(0)
    
    return Response(
        content=buffer.getvalue(),
        media_type="application/pdf",
        headers={"Content-Disposition": f"attachment; filename={request.ticker}_tearsheet.pdf"}
    )


def generate_ppt_tearsheet(request: TearsheetRequest):
    """Generate PowerPoint tearsheet using python-pptx."""
    prs = Presentation()
    
    # Title slide
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    
    title.text = f"{request.ticker} - Investment Tearsheet"
    subtitle.text = f"Generated: {datetime.now().strftime('%B %d, %Y')}\nDealLens - Investment Banking Platform"
    
    # Add a slide for each section
    for section in request.sections:
        bullet_slide_layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(bullet_slide_layout)
        shapes = slide.shapes
        
        title_shape = shapes.title
        body_shape = shapes.placeholders[1]
        
        title_shape.text = f"{section.title()} Analysis"
        
        tf = body_shape.text_frame
        tf.text = f"Key insights for {request.ticker} {section} analysis:"
        
        # Add some mock bullet points
        for i in range(3):
            p = tf.add_paragraph()
            p.text = f"Analysis point {i+1} for {section} section"
            p.level = 1
    
    # Save to buffer
    buffer = io.BytesIO()
    prs.save(buffer)
    buffer.seek(0)
    
    return Response(
        content=buffer.getvalue(),
        media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
        headers={"Content-Disposition": f"attachment; filename={request.ticker}_tearsheet.pptx"}
    )
